# -*- coding: utf-8 -*-
"""KAASA smartfarmingsight 종합 도입 제안서 (PPTX, 16:9)
   = 통합소개서(시스템 소개 + 7단계 로드맵) 기반 + 제안서 요소(문제·차별점·무료 요금제·회원가입 CTA)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

SHOTS = os.path.join("out", "ppt_shots")
OUT = os.path.join("out", "KAASA_smartfarmingsight_제안서.pptx")

GREEN_DARK = RGBColor(0x0F,0x51,0x32); GREEN = RGBColor(0x2E,0xCC,0x71)
ORANGE = RGBColor(0xE6,0x7E,0x22); INK = RGBColor(0x1A,0x1A,0x1A)
GRAY = RGBColor(0x5A,0x6A,0x60); LIGHT = RGBColor(0xF2,0xF6,0xF4)
RED = RGBColor(0xC0,0x39,0x2B); BLUE = RGBColor(0x2D,0x9C,0xDB)
PURPLE = RGBColor(0x8E,0x44,0xAD); MINT = RGBColor(0xCF,0xE9,0xDB); WHITE = RGBColor(0xFF,0xFF,0xFF)
FONT = "맑은 고딕"

prs = Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]; PR=390/844

def _box(s,x,y,w,h,fill=None,line=None,lw=1):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    return sp
def _round(s,x,y,w,h,fill,line=None):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h); sp.shadow.inherit=False
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1.5)
    return sp
def _text(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sa=6):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sa); p.space_before=Pt(0)
        if isinstance(para,tuple): para=[para]
        for (t,sz,b,c) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=FONT
    return tb
def _bar(s,title,kicker=None):
    _box(s,0,0,SW,Inches(0.95),fill=GREEN_DARK); _box(s,0,Inches(0.95),SW,Inches(0.06),fill=GREEN)
    _text(s,Inches(0.6),0,Inches(10.5),Inches(0.95),[[(title,24,True,WHITE)]],anchor=MSO_ANCHOR.MIDDLE)
    if kicker:_text(s,Inches(9.3),0,Inches(3.4),Inches(0.95),[[(kicker,12,False,MINT)]],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
def phone(s,name,cx,cy,h_in):
    h=Inches(h_in); w=Emu(int(h*PR)); x=cx-w//2; pad=Inches(0.05)
    _box(s,x-pad,cy-pad,w+pad*2,h+pad*2,fill=WHITE,line=RGBColor(0xD8,0xE2,0xDC))
    p=os.path.join(SHOTS,name+".png")
    if os.path.exists(p): s.shapes.add_picture(p,x,cy,height=h)
def caption(s,t,cx,y,w=2.6):
    _text(s,cx-Inches(w/2),y,Inches(w),Inches(0.4),[[(t,12,True,GREEN_DARK)]],align=PP_ALIGN.CENTER)
def bullets(s,x,y,w,h,items,size=14,head=None):
    runs=[[(head,17,True,GREEN_DARK)]] if head else []
    runs+=[[("•  ",size,True,GREEN),(it,size,False,INK)] for it in items]
    _text(s,x,y,w,h,runs,sa=9)
def cards4(s,items,y0=Inches(1.45)):
    cw,ch,gx,gy=Inches(5.85),Inches(2.15),Inches(0.3),Inches(0.3); x0=Inches(0.6)
    for i,(ico,t,d,col) in enumerate(items):
        r,c=divmod(i,2); x=x0+c*(cw+gx); y=y0+r*(ch+gy)
        _round(s,x,y,cw,ch,LIGHT,line=RGBColor(0xDD,0xE7,0xE1)); _box(s,x,y,Inches(0.14),ch,fill=col)
        _text(s,x+Inches(0.35),y+Inches(0.22),cw-Inches(0.6),Inches(0.6),[[(ico+"  "+t,17,True,col)]])
        _text(s,x+Inches(0.35),y+Inches(0.92),cw-Inches(0.6),ch-Inches(1.1),[[(d,13.5,False,INK)]])
def feat(title,kicker,head,items,shots):
    s=prs.slides.add_slide(BLANK); _bar(s,title,kicker)
    bullets(s,Inches(0.6),Inches(1.45),Inches(5.1),Inches(5.4),items,size=14,head=head)
    for i,(nm,cap) in enumerate(shots[:2]):
        cx=Inches(8.0) if i==0 else Inches(11.0); phone(s,nm,cx,Inches(1.55),5.0); caption(s,cap,cx,Inches(6.95))
def divider(title,subtitle):
    s=prs.slides.add_slide(BLANK); _box(s,0,0,SW,SH,fill=GREEN_DARK)
    _box(s,Inches(0.9),Inches(3.0),Inches(0.18),Inches(1.5),fill=GREEN)
    _text(s,Inches(1.3),Inches(3.0),Inches(11),Inches(1.0),[[(title,36,True,WHITE)]],anchor=MSO_ANCHOR.MIDDLE)
    _text(s,Inches(1.32),Inches(4.15),Inches(11),Inches(0.6),[[(subtitle,16,False,MINT)]])

# ═════ 1. 표지(제안서) ═════
s=prs.slides.add_slide(BLANK); _box(s,0,0,SW,SH,fill=GREEN_DARK)
if os.path.exists("og-image.png"):
    bw=Inches(7.4); s.shapes.add_picture("og-image.png",(SW-bw)//2,Inches(0.55),width=bw)
_text(s,Inches(1),Inches(4.55),Inches(11.3),Inches(1.5),
      [[("스마트팜 경영최적화 도입 제안서",32,True,WHITE)],
       [("환경관리 · 생육모델 · 이기종 통합 · 온실/노지 — 데이터로 농사를 결정하다",15,False,MINT)]],align=PP_ALIGN.CENTER)
_round(s,Inches(4.85),Inches(6.35),Inches(3.63),Inches(0.6),GREEN)
_text(s,Inches(4.85),Inches(6.35),Inches(3.63),Inches(0.6),[[("무료로 시작하기 →  farmingsight.org",13,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# ═════ 2. 문제(페인) ═════
s=prs.slides.add_slide(BLANK); _bar(s,"이런 고민, 있으셨나요?","Problem")
cards4(s,[("🤔","감과 경험에 의존","환경·관수·시비 판단이 데이터 없이 직관에 의존 → 작황 편차가 크고 재현이 어렵습니다.",RED),
          ("🧩","데이터가 흩어져 있음","센서·기상·시장·장비 데이터가 단절돼 경영 의사결정에 활용하지 못합니다.",ORANGE),
          ("🔌","이기종 장비 파편화","제조사마다 제어기·앱이 달라 통합 관제·비교가 불가능합니다.",RED),
          ("💸","경영이 깜깜이","수확량·매출·비용·순이익을 미리 못 봐 투자·출하 판단이 어렵습니다.",ORANGE)])

# ═════ 3. 해결(가치 제안) ═════
s=prs.slides.add_slide(BLANK); _bar(s,"KAASA가 해결해 드립니다","Solution")
cards4(s,[("📊","데이터 기반 의사결정","환경·관수·시비를 일사적산·VPD·전략표로 정밀 처방 — 감 대신 데이터로.",GREEN_DARK),
          ("🗂","흩어진 데이터를 하나로","센서·기상·시장·장비를 표준 변수로 통합해 한 화면에서 운영.",GREEN_DARK),
          ("🔗","이기종 장비 표준 통합","서로 다른 제조사 제어기·센서를 MQTT·게이트웨이로 단일 관제.",GREEN_DARK),
          ("📈","수확량·순이익 예측 경영","M1~M5 모델로 수확량→매출→비용→순이익을 미리 예측·최적화.",GREEN_DARK)])

# ═════ 4. 시스템 개요 ═════
s=prs.slides.add_slide(BLANK); _bar(s,"시스템 개요","Overview")
_text(s,Inches(0.6),Inches(1.25),Inches(12.1),Inches(0.9),
      [[("환경·기상·시장 데이터를 받아 ",16,False,INK),("수확량 → 매출 → 비용 → 순이익",16,True,GREEN_DARK),
        ("을 예측하고, 최적 환경과 관수·출하 의사결정을 제시하는 스마트팜 운영 OS",16,False,INK)]])
ov=[("🌡 스마트팜 환경관리","온·습도·CO₂·VPD·DLI 전략표·AI 처방"),("💧 일사적산 관수","P1~P6 일일 곡선·EC·배액률·함수율"),
    ("📈 생육모델·수확량 예측","M1~M5 모델 체인 예측"),("🔌 이기종 장비 통합","제조사 무관 단일 관제"),
    ("🏡 온실6+노지7 = 13작목","광역 작목 지원"),("🤝 공동출하·경영분석","시세 비교·월간 리포트")]
cx0,cy0,cw,ch,gx,gy=Inches(0.6),Inches(2.25),Inches(3.9),Inches(1.45),Inches(0.18),Inches(0.22)
for i,(t,d) in enumerate(ov):
    r,c=divmod(i,3); x=cx0+c*(cw+gx); y=cy0+r*(ch+gy)
    _box(s,x,y,cw,ch,fill=LIGHT,line=RGBColor(0xDD,0xE7,0xE1))
    _text(s,x+Inches(0.18),y+Inches(0.12),cw-Inches(0.36),ch-Inches(0.24),[[(t,14,True,GREEN_DARK)],[(d,11.5,False,GRAY)]],sa=4)

# ═════ 5. 아키텍처 ═════
s=prs.slides.add_slide(BLANK); _bar(s,"핵심 아키텍처","Architecture")
layers=[("프런트엔드","모바일 화면 41종 · 공용 레이어 · PWA(오프라인·캐시)",GREEN),
        ("API (FastAPI)","JWT 인증 · 농가 소유권 · 라우터(farmer·admin·billing·ws…)",BLUE),
        ("분석 코어","M1 생육→M2 수확량→M3 매출→M4 비용→M5 병해 · profit_optimizer",ORANGE),
        ("파이프라인","ETL 증분 · 임계 재학습·배포 게이트 · 레지스트리 · 연합학습 · MQTT",PURPLE),
        ("배포·운영","Cloudflare tunnel→uvicorn:8000 · watchdog 자동복구 · farmingsight.org",GREEN_DARK)]
y=Inches(1.4)
for t,d,col in layers:
    _box(s,Inches(0.6),y,Inches(12.1),Inches(0.98),fill=LIGHT,line=RGBColor(0xDD,0xE7,0xE1)); _box(s,Inches(0.6),y,Inches(0.16),Inches(0.98),fill=col)
    _text(s,Inches(1.0),y,Inches(3.0),Inches(0.98),[[(t,16,True,col)]],anchor=MSO_ANCHOR.MIDDLE)
    _text(s,Inches(4.0),y,Inches(8.5),Inches(0.98),[[(d,12.5,False,INK)]],anchor=MSO_ANCHOR.MIDDLE); y+=Inches(1.12)

# ═════ 6~13. 기능 화면 8 ═════
feat("랜딩 · 화면 네비게이터","Feature","첫 진입 & 탐색",["키워드 중심 랜딩(SEO)으로 가치 전달","41개 화면 네비게이터","등급별 맞춤 메뉴·잠금 배지","PWA 설치·오프라인 캐시"],[("intro","랜딩"),("smartos","화면 네비게이터")])
feat("통합 홈 · 온실 홈","Feature","일일 운영 대시보드",["VPD·관수·생육·병해 핵심 지표","의사결정 카드 + 적용 기록 폐루프","실시간 센서 WebSocket","작물 전환 13작목 비교"],[("c3_home","통합 홈"),("g1_home","온실 홈")])
feat("환경관리 전략 · 관수 P1~P6","Feature","정밀 제어",["생육시기×4구간 전략표","편차→AI 처방(제어 4단계)","일사 적산(J/cm²) P1~P6","EC·pH·배액률·함수율"],[("g2_env","환경관리 전략"),("g3_period","관수 P1~P6")])
feat("생육·수확량 예측 · 병해 경보","Feature","예측과 예방",["작물별 수확량(kg/m²) 예측","예측 vs 실측 → 재학습","결로시간 기반 병해 조기경보","방제 이행 기록 폐루프"],[("g4_growth","수확량 예측"),("g5_disease","병해 조기경보")])
feat("노지 관리 · 경영·ERP","Feature","노지와 경영",["제주 노지 7작목·토양수분·관개","손익(매출·비용·순이익) 분해","절감 조치 기록 비용 최적화","현장 입력→모델 학습"],[("f1_field","노지 관리"),("c5_erp","경영·ERP 분석")])
feat("공동출하 · 경영성과 리포트","Feature","유통과 성과",["시세·채널 비교·공동출하 참여","작목별 등급 자동 적용","월간 경영성과 리포트","성과지표 변화율·월 스냅샷"],[("c12_joint","공동출하"),("c14_report","경영성과 리포트")])
feat("등급 비교 · 클러스터 관제","Feature","구독과 광역 관제",["4등급 기능 매트릭스·AI 쿼터","업그레이드 안내","다중농가 클러스터 집계","이상 농가 위치특정(PII 익명화)"],[("c22_tiers","등급 비교"),("c20_cluster","클러스터 관제")])
feat("AI 영농비서 · 이기종 통합","Feature","대화형 운영과 통합",["내 농장 데이터 AI 비서(RAG)","관수·환경·수확·병해 질의응답","제조사 무관 제어기·센서 등록","연동 신청으로 단일 화면"],[("c13_chat","AI 영농비서"),("c8_interop","이기종 장비 통합")])

# ═════ 14. 작물·모델 성능 ═════
s=prs.slides.add_slide(BLANK); _bar(s,"작물 커버리지 & 모델 성능","Crops & Models")
bullets(s,Inches(0.6),Inches(1.35),Inches(6.2),Inches(5.4),
  ["온실 6작목 + 제주 노지 7작목 = 총 13작목","M1 생육·M2 수확량·M3 매출·M4 비용·M5 병해",
   "수확량 예측 MAPE 8.7~19.8%(게이트 35% 이내)","생육 설명력 R² 0.17~0.37 — 데이터 축적 시 개선",
   "드리프트 모니터링 + 임계 자동 재학습·폴백"],size=14.5,head="13작목 × M1~M5 모델 체인")
phone(s,"g4_growth",Inches(10.2),Inches(1.55),5.2); caption(s,"생육·수확량 예측",Inches(10.2),Inches(7.0))

# ═════ 15. 왜 KAASA인가 ═════
s=prs.slides.add_slide(BLANK); _bar(s,"왜 KAASA인가","Why us")
diffs=[("🌞 일사적산 정밀 관수","시각이 아닌 광량(J/cm²) 기반 — Priva·Grodan 글로벌 기준 정합"),
       ("🧪 13작목 검증 모델","온실 6+노지 7, 수확량 예측 오차(MAPE) 8.7~19.8%"),
       ("🔌 이기종 장비 통합","제조사 무관 표준 변수 매핑·MQTT·게이트웨이 단일 관제"),
       ("🔁 폐루프 학습","현장 이행 데이터가 모델을 지속 개선 — 쓸수록 똑똑해짐"),
       ("🛡 운영급 보안·안정","외부 보안감사 전 항목 조치·HTTPS·PWA·자동복구")]
y=Inches(1.4)
for t,d in diffs:
    _round(s,Inches(0.6),y,Inches(12.1),Inches(0.92),LIGHT,line=RGBColor(0xDD,0xE7,0xE1))
    _text(s,Inches(0.95),y,Inches(4.2),Inches(0.92),[[(t,15,True,GREEN_DARK)]],anchor=MSO_ANCHOR.MIDDLE)
    _text(s,Inches(5.2),y,Inches(7.3),Inches(0.92),[[(d,13,False,INK)]],anchor=MSO_ANCHOR.MIDDLE); y+=Inches(1.05)

# ═════ 16. 보안·운영 ═════
s=prs.slides.add_slide(BLANK); _bar(s,"보안 · 운영","Security & Ops")
bullets(s,Inches(0.6),Inches(1.4),Inches(12.1),Inches(2.6),
  ["무인증 데이터 수집·모델오염·DoS 차단(인증+소유권)","WebSocket/센서 무인증 도청 차단",
   "공개 데모 무자격 토큰 — 소스 내 자격증명 제거","fail-closed 인증·농가 격리·결제/연합 무결성·클러스터 PII 익명화"],
  size=14.5,head="외부 보안감사 P0~P3 전 항목 조치 완료")
bullets(s,Inches(0.6),Inches(4.4),Inches(12.1),Inches(2.6),
  ["Cloudflare tunnel → uvicorn :8000 (HTTPS)","watchdog 30초 자동복구 + 단일 인스턴스 가드",
   "PUBLIC_DEMO 읽기전용 게이트","PWA 캐시 버저닝 · SEO(robots·sitemap·구조화데이터)"],
  size=14.5,head="운영 인프라")

# ═════ 17~25. 7단계 로드맵 ═════
divider("7단계 실증 추진 로드맵","참여농가 사전진단부터 데이터 수집·AI 의사결정·컨설팅·사업화 모델 도출까지")
s=prs.slides.add_slide(BLANK); _bar(s,"실증 추진 7단계 — 한눈에","Roadmap")
ss=[("1","참여농가 사전진단 & 현장 문제 유형화","C17·C18"),("2","데이터 수집항목·주기·입력방식 확정","C21·C2"),
    ("3","농가별 데이터 수집 & 품질관리","C16·C8"),("4","AI 의사결정 지원 대시보드 시범적용","C3·G2"),
    ("5","농가별 분석리포트 & 현장 컨설팅","C14·C4"),("6","농가 피드백 수렴 & 서비스 개선","C13·도움말"),
    ("7","실증성과 분석 & 사업화 모델 도출","C9·C10")]
y=Inches(1.3)
for n,t,sc in ss:
    _box(s,Inches(0.6),y,Inches(0.62),Inches(0.62),fill=GREEN); _text(s,Inches(0.6),y,Inches(0.62),Inches(0.62),[[(n,20,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    _box(s,Inches(1.4),y,Inches(9.0),Inches(0.62),fill=LIGHT,line=RGBColor(0xDD,0xE7,0xE1)); _text(s,Inches(1.65),y,Inches(8.6),Inches(0.62),[[(t,14,True,INK)]],anchor=MSO_ANCHOR.MIDDLE)
    _box(s,Inches(10.6),y,Inches(2.1),Inches(0.62),fill=GREEN_DARK); _text(s,Inches(10.6),y,Inches(2.1),Inches(0.62),[[(sc,12,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(0.78)
ROAD=[("1단계 · 사전진단 & 현장 문제 유형화","Step 1",["[C17] 성숙도 점수·6대영역 진단·우선순위 처방·ROI·결과 PDF","[C18] RDA 현장 체크리스트로 센서 밖 항목까지 반영해 문제 유형화"],[("c17_diagnosis","C17 시스템 종합진단"),("c18_checklist","C18 현장 컨설팅 문진")]),
      ("2단계 · 수집항목·주기·입력방식 확정","Step 2",["[C21] 장비·서비스·프로토콜·인프라정도로 이기종·외부데이터·전문가 연동 신청·관리","[C2] 활용 모드·항목별 공개 토글로 수집·활용범위 확정"],[("c21_apply","C21 연동·서비스 신청"),("c2_consent","C2 데이터 동의")]),
      ("3단계 · 데이터 수집 & 품질관리","Step 3",["[C16] 이기종 장비 등록·표준변수 매핑(견적서·시방서·PDF 자동추출)","[C8] 연동API·연동률·MQTT·게이트웨이로 품질·연동상태 관리"],[("c16_equipment","C16 시설기자재 등록"),("c8_interop","C8 이기종 연동")]),
      ("4단계 · AI 의사결정 대시보드 시범적용","Step 4",["[C3] 오늘의 결정·수확예측·예상매출·VPD·AI 추천 통합 표시","[G2] 실측·전략표·제어모드 4단계·AI 에이전트(30초)로 시범적용"],[("c3_home","C3 통합 홈"),("g2_env","G2 환경·기후 제어")]),
      ("5단계 · 분석리포트 & 현장 컨설팅","Step 5",["[C14] 6대영역·9대 성과지표·이행활동·실행 To-do·PDF","[C4] AI 종합진단·개선추천·환경이상 감지·전문가 컨설팅 예약"],[("c14_report","C14 월간 경영성과 리포트"),("c4_diagnosis","C4 AI 진단")]),
      ("6단계 · 피드백 수렴 & 기능 개선","Step 6",["[C13] 농장데이터+농진청·병해충DB(RAG) AI 챗봇 Q&A·피드백 수렴","[도움말] 가이드·FAQ·용어·출처배지 매뉴얼로 사용성 개선"],[("c13_chat","C13 AI 영농비서"),("help","도움말·사용자 매뉴얼")]),
      ("7단계 · 실증성과 분석 & 사업화","Step 7",["[C9] RDA 표준·우수농가 대비 비교·개선포인트·월간리포트 연계","[C10] 연매출·연에너지비 기반 ROI·투자안 비교·회수기간으로 사업화"],[("c9_benchmark","C9 벤치마킹"),("c10_roi","C10 투자 ROI")])]
for title,kicker,items,shots in ROAD: feat(title,kicker,"추진 내용",items,shots)

# ═════ 26. 무료 요금제 ═════
s=prs.slides.add_slide(BLANK); _bar(s,"무료로 시작하세요 · 요금제","Pricing")
tiers=[("Basic","무료","핵심 대시보드·기록·기본 분석",GREEN),("Smart","구독","환경 전략표·관수 정밀·AI 추천",BLUE),
       ("Pro","구독","수확 예측·경영 리포트·벤치마크",ORANGE),("Enterprise","구독","클러스터 관제·연동·전문 컨설팅",PURPLE)]
cw,gx=Inches(2.85),Inches(0.27); x=Inches(0.7)
for nm,price,desc,col in tiers:
    _round(s,x,Inches(1.5),cw,Inches(3.6),WHITE,line=col); _box(s,x,Inches(1.5),cw,Inches(0.85),fill=col)
    _text(s,x,Inches(1.5),cw,Inches(0.85),[[(nm,17,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    _text(s,x,Inches(2.5),cw,Inches(0.7),[[(price,22,True,col)]],align=PP_ALIGN.CENTER)
    _text(s,x+Inches(0.25),Inches(3.3),cw-Inches(0.5),Inches(1.6),[[(desc,12.5,False,INK)]],align=PP_ALIGN.CENTER); x+=cw+gx
_text(s,Inches(0.7),Inches(5.4),Inches(12),Inches(0.6),[[("✓ Basic 무료 플랜으로 즉시 시작 · 신용카드 불필요 · 언제든 업그레이드",15,True,GREEN_DARK)]],align=PP_ALIGN.CENTER)
phone(s,"c22_tiers",Inches(11.4),Inches(5.95),1.35)

# ═════ 27. ★ 회원가입 CTA ═════
s=prs.slides.add_slide(BLANK); _box(s,0,0,SW,SH,fill=GREEN_DARK)
_text(s,Inches(0.8),Inches(0.7),Inches(8.0),Inches(1.2),
      [[("지금, 무료로 시작하세요",34,True,WHITE)],[("회원가입 후 농장을 등록하면 AI 진단을 무료로 체험할 수 있습니다.",16,False,MINT)]])
cta=[("1","farmingsight.org 접속","웹·모바일 어디서나, 설치 없이"),("2","회원가입(이메일·전화)","역할 선택 — 농가·조합·유통·전문가·공공"),
     ("3","농장 세팅","지역·작목·재배방식·장비 등록"),("4","AI 진단 무료 체험","시스템 종합진단·환경/관수 처방 즉시 확인")]
y=Inches(2.1)
for n,t,d in cta:
    _box(s,Inches(0.8),y,Inches(0.55),Inches(0.55),fill=GREEN); _text(s,Inches(0.8),y,Inches(0.55),Inches(0.55),[[(n,17,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    _text(s,Inches(1.55),y-Inches(0.05),Inches(6.6),Inches(0.8),[[(t,16,True,WHITE)],[(d,12,False,MINT)]],sa=0); y+=Inches(0.95)
_box(s,Inches(8.7),Inches(1.9),Inches(3.9),Inches(4.7),fill=WHITE)
if os.path.exists("out/qr_farmingsight.png"): s.shapes.add_picture("out/qr_farmingsight.png",Inches(9.45),Inches(2.2),height=Inches(2.3))
_text(s,Inches(8.7),Inches(4.65),Inches(3.9),Inches(0.4),[[("QR 스캔 → 바로 접속",13,True,GREEN_DARK)]],align=PP_ALIGN.CENTER)
_text(s,Inches(8.7),Inches(5.05),Inches(3.9),Inches(0.4),[[("farmingsight.org",17,True,GREEN_DARK)]],align=PP_ALIGN.CENTER)
_text(s,Inches(8.7),Inches(5.6),Inches(3.9),Inches(0.9),[[("✓ 무료 플랜 · 신용카드 불필요",12.5,True,GREEN_DARK)],[("✓ 13작목 · 온실/노지 지원",12.5,True,GREEN_DARK)]],align=PP_ALIGN.CENTER,sa=3)

# ═════ 28. 마무리 ═════
s=prs.slides.add_slide(BLANK); _box(s,0,0,SW,SH,fill=GREEN_DARK)
_text(s,Inches(1),Inches(2.6),Inches(11.3),Inches(2.0),[[("함께 시작하시죠",38,True,WHITE)],[("KAASA smartfarmingsight · 데이터로 농사를 결정하다",17,False,MINT)]],align=PP_ALIGN.CENTER)
_round(s,Inches(4.7),Inches(4.7),Inches(3.93),Inches(0.7),GREEN)
_text(s,Inches(4.7),Inches(4.7),Inches(3.93),Inches(0.7),[[("무료 회원가입 →  farmingsight.org",14,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT); print("저장 완료:", OUT, "| 슬라이드", len(prs.slides._sldIdLst))
