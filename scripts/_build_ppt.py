# -*- coding: utf-8 -*-
"""KAASA Farmingsight 시스템 소개서 (PPTX, 16:9) — 화면 스크린샷 포함."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

SHOTS = os.path.join("out", "ppt_shots")
OUT = os.path.join("out", "KAASA_Farmingsight_소개서.pptx")

# 팔레트
GREEN_DARK = RGBColor(0x0F, 0x51, 0x32)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
INK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x5A, 0x6A, 0x60)
LIGHT = RGBColor(0xF2, 0xF6, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

PHONE_RATIO = 390 / 844  # w/h


def _box(slide, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    return sp


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=6, wrap=True):
    """runs: list of (text, size, bold, color) per paragraph."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        if isinstance(para, tuple): para = [para]
        for (t, sz, bold, col) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col
            r.font.name = "맑은 고딕"
    return tb


def _bar(slide, title, kicker=None):
    _box(slide, 0, 0, SW, Inches(0.95), fill=GREEN_DARK)
    _box(slide, 0, Inches(0.95), SW, Inches(0.06), fill=GREEN)
    _text(slide, Inches(0.6), 0, Inches(11), Inches(0.95),
          [[(title, 26, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        _text(slide, Inches(9.3), 0, Inches(3.4), Inches(0.95),
              [[(kicker, 12, False, RGBColor(0xCF,0xE9,0xDB))]],
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def phone(slide, name, cx, cy, h_in):
    """폰 스크린샷을 (cx,cy 중심 상단) 배치. 반환 width."""
    h = Inches(h_in); w = Emu(int(h * PHONE_RATIO))
    path = os.path.join(SHOTS, name + ".png")
    x = cx - w // 2
    # 테두리 카드
    pad = Inches(0.05)
    _box(slide, x - pad, cy - pad, w + pad*2, h + pad*2, fill=WHITE, line=RGBColor(0xD8,0xE2,0xDC))
    slide.shapes.add_picture(path, x, cy, height=h)
    return w


def caption(slide, text, cx, y, w_in=3.0):
    _text(slide, cx - Inches(w_in/2), y, Inches(w_in), Inches(0.4),
          [[(text, 12, True, GREEN_DARK)]], align=PP_ALIGN.CENTER)


def bullets(slide, x, y, w, h, items, size=14, head=None):
    runs = []
    if head: runs.append([(head, 17, True, GREEN_DARK)])
    for it in items:
        runs.append([("•  ", size, True, GREEN), (it, size, False, INK)])
    _text(slide, x, y, w, h, runs, space_after=9)


# ───────────────────────── 1. 타이틀 ─────────────────────────
s = prs.slides.add_slide(BLANK)
_box(s, 0, 0, SW, SH, fill=GREEN_DARK)
if os.path.exists("og-image.png"):
    # 상단 배너(1200x630)
    bw = Inches(8.2); bh = Emu(int(bw * 630/1200))
    s.shapes.add_picture("og-image.png", (SW-bw)//2, Inches(0.7), width=bw)
_text(s, Inches(1), Inches(5.0), Inches(11.3), Inches(1.2),
      [[("스마트농업 경영최적화 의사결정 플랫폼", 30, True, WHITE)],
       [("환경관리 · 생육모델 · 이기종 통합 · 온실/노지 관리 — 시스템 소개서", 16, False, RGBColor(0xCF,0xE9,0xDB))]],
      align=PP_ALIGN.CENTER)
_text(s, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5),
      [[("https://farmingsight.org", 13, False, GREEN)]], align=PP_ALIGN.CENTER)

# ───────────────────────── 2. 시스템 개요 ─────────────────────────
s = prs.slides.add_slide(BLANK); _bar(s, "시스템 개요", "What it is")
_text(s, Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.9),
      [[("환경·기상·시장 데이터를 받아 ", 16, False, INK),
        ("수확량 → 매출 → 비용 → 순이익", 16, True, GREEN_DARK),
        ("을 예측하고, 최적 환경과 관수·출하 의사결정을 제시하는 스마트팜 운영 OS", 16, False, INK)]])
cards = [
    ("🌡 스마트팜 환경관리", "온도·습도·CO₂·VPD·DLI 광연동 전략표와 AI 처방"),
    ("💧 일사적산 관수", "P1~P6 일일 관수 곡선(J/cm²)·EC·배액률·함수율"),
    ("📈 생육모델·수확량 예측", "M1~M5 모델 체인으로 작물별 수확량·순이익 예측"),
    ("🔌 이기종 장비 통합", "서로 다른 제조사 복합환경제어기·센서를 하나로"),
    ("🏡 온실 6 + 노지 7 = 13작목", "온실 6종·제주 노지 7종 광역 지원"),
    ("🤝 공동출하·경영분석", "시세·채널 비교, 월간 경영성과 리포트"),
]
cx0, cy0, cw, ch, gx, gy = Inches(0.6), Inches(2.25), Inches(3.9), Inches(1.45), Inches(0.18), Inches(0.22)
for i, (t, d) in enumerate(cards):
    r, c = divmod(i, 3)
    x = cx0 + c*(cw+gx); y = cy0 + r*(ch+gy)
    _box(s, x, y, cw, ch, fill=LIGHT, line=RGBColor(0xDD,0xE7,0xE1))
    _text(s, x+Inches(0.18), y+Inches(0.12), cw-Inches(0.36), ch-Inches(0.24),
          [[(t, 14.5, True, GREEN_DARK)], [(d, 11.5, False, GRAY)]], space_after=4)

# ───────────────────────── 3. 아키텍처 ─────────────────────────
s = prs.slides.add_slide(BLANK); _bar(s, "핵심 아키텍처", "Architecture")
layers = [
    ("프런트엔드", "모바일 화면 41종(screens/) · 공용 레이어(components/) · PWA(SW 캐시·오프라인)", GREEN),
    ("API (FastAPI)", "JWT 인증 · 농가 소유권 · PUBLIC_DEMO 게이트 · 라우터(farmer·admin·auth·billing·federated·ws)", RGBColor(0x2D,0x9C,0xDB)),
    ("분석 코어", "M1 생육 → M2 수확량 → M3 매출 → M4 비용 → M5 병해 · profit_optimizer 추천", RGBColor(0xE6,0x7E,0x22)),
    ("파이프라인", "ETL 증분 병합 · 임계 재학습·배포 게이트 · 모델 레지스트리 · 연합학습 · MQTT 수집", RGBColor(0x8E,0x44,0xAD)),
    ("배포·운영", "Cloudflare named tunnel → uvicorn :8000 · watchdog 자동복구 · farmingsight.org", GREEN_DARK),
]
y = Inches(1.4)
for title, desc, col in layers:
    _box(s, Inches(0.6), y, Inches(12.1), Inches(0.98), fill=LIGHT, line=RGBColor(0xDD,0xE7,0xE1))
    _box(s, Inches(0.6), y, Inches(0.16), Inches(0.98), fill=col)
    _text(s, Inches(1.0), y, Inches(3.0), Inches(0.98), [[(title, 16, True, col)]], anchor=MSO_ANCHOR.MIDDLE)
    _text(s, Inches(4.0), y, Inches(8.5), Inches(0.98), [[(desc, 12.5, False, INK)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.12)

# ───────────────────────── 화면 슬라이드 빌더 ─────────────────────────
def feature_slide(title, kicker, left_head, left_items, shots):
    s = prs.slides.add_slide(BLANK); _bar(s, title, kicker)
    bullets(s, Inches(0.6), Inches(1.45), Inches(5.1), Inches(5.4), left_items, size=14.5, head=left_head)
    # 우측 폰 1~2개
    area_cx = Inches(9.4)
    if len(shots) == 1:
        w = phone(s, shots[0][0], area_cx, Inches(1.5), 5.4)
        caption(s, shots[0][1], area_cx, Inches(7.0))
    else:
        gap = Inches(1.9)
        for i, (nm, cap) in enumerate(shots[:2]):
            cx = area_cx - gap//1 + i*gap if False else (Inches(8.0) if i==0 else Inches(11.0))
            phone(s, nm, cx, Inches(1.55), 5.0)
            caption(s, cap, cx, Inches(6.95), w_in=2.6)
    return s

feature_slide("랜딩 · 화면 네비게이터", "Onboarding",
    "첫 진입 & 탐색",
    ["키워드 중심 랜딩(SEO 최적화)으로 가치 전달",
     "화면 네비게이터에서 41개 화면을 한눈에 탐색",
     "등급별 맞춤 메뉴·잠금 배지로 접근 경로 안내",
     "PWA 설치·오프라인 캐시 지원"],
    [("intro","랜딩(소개)"), ("smartos","화면 네비게이터")])

feature_slide("통합 홈 · 온실 홈", "Dashboard",
    "일일 운영 대시보드",
    ["VPD·관수·생육·병해 핵심 지표를 한 화면에",
     "DecisionDeck 의사결정 카드 + '적용 기록' 폐루프",
     "실시간 센서 WebSocket 연동(인증·소유권 검증)",
     "작물(농장) 전환으로 13작목 즉시 비교"],
    [("c3_home","통합 홈"), ("g1_home","온실 홈")])

feature_slide("환경관리 전략 · 관수 P1~P6", "Climate & Irrigation",
    "정밀 제어",
    ["생육시기 × 하루 4구간 환경 전략표(온·습도·CO₂·VPD)",
     "목표 대비 실측 편차 → AI 처방 자동 생성",
     "관수는 시각이 아닌 일사 적산(J/cm²) 기준 P1~P6",
     "급액/배액 EC·pH·배액률·함수율 dry-back 관리"],
    [("g2_env","환경관리 전략"), ("g3_period","관수 P1~P6")])

feature_slide("생육·수확량 예측 · 병해 조기경보", "ML & IPM",
    "예측과 예방",
    ["M1~M2 생육모델로 작물별 수확량(kg/m²) 예측",
     "예측 vs 실측 피드백 → 임계 도달 시 재학습",
     "결로시간 기반 병해 감염위험 조기경보(IPM)",
     "방제 이행 기록으로 학습·품질 폐루프 완성"],
    [("g4_growth","생육·수확량 예측"), ("g5_disease","병해 조기경보")])

feature_slide("노지 관리 · 경영·ERP 분석", "Field & ERP",
    "노지와 경영",
    ["제주 노지 7작목 — 토양수분·관개·기상 기반 운영",
     "경영 손익(매출·비용·순이익) 분해 분석",
     "절감 조치 이행 기록으로 비용 최적화 루프",
     "현장 데이터 입력 → 모델 학습 피드"],
    [("f1_field","노지 관리"), ("c5_erp","경영·ERP 분석")])

feature_slide("공동출하 · 월간 경영성과 리포트", "Market & Report",
    "유통과 성과",
    ["시세·유통 채널 비교, 공동출하 참여 신청",
     "작목별 등급 체계 자동 적용(번과/특·상·등외 등)",
     "월간 경영성과 리포트(정책 지표 연동)",
     "성과지표 변화율·월 스냅샷 이력 추적"],
    [("c12_joint","공동출하"), ("c14_report","월간 경영성과 리포트")])

feature_slide("등급 비교 · 클러스터 관제", "Tiers & Cluster",
    "구독과 광역 관제",
    ["4등급(basic·smart·pro·enterprise) 기능 매트릭스",
     "등급별 접근 경로·AI 쿼터 가시화 + 업그레이드",
     "다중농가 클러스터 광역 작황·진단 집계",
     "이상 농가 위치특정(공개 응답은 PII 익명화)"],
    [("c22_tiers","등급 비교"), ("c20_cluster","클러스터 관제")])

feature_slide("AI 영농비서 · 이기종 장비 통합", "AI & Interop",
    "대화형 운영과 통합",
    ["내 농장 데이터 기반 AI 영농비서(규칙·LLM 폴백)",
     "관수·환경·평년기상·수확·병해 질의응답",
     "서로 다른 제조사 복합환경제어기·센서 등록",
     "데이터 통합·연동 신청으로 단일 운영 화면"],
    [("c13_chat","AI 영농비서"), ("c8_interop","이기종 장비 통합")])

# ───────────────────────── 작물·모델 ─────────────────────────
s = prs.slides.add_slide(BLANK); _bar(s, "작물 커버리지 & 모델", "Crops & Models")
bullets(s, Inches(0.6), Inches(1.4), Inches(6.0), Inches(5.4),
    ["온실 6작목: 딸기·오이·완숙토마토·방울토마토·파프리카·참외",
     "제주 노지 7작목 — 총 13작목 지원",
     "M1 생육 · M2 수확량 · M3 매출 · M4 비용 · M5 병해",
     "배포 게이트(R²/MAPE/F1)로 신뢰도 검증",
     "게이트 미달 작물은 통계 폴백으로 안전 동작",
     "드리프트 모니터링 + 임계 기반 자동 재학습"],
    size=15, head="13작목 × M1~M5 모델 체인")
phone(s, "g4_growth", Inches(10.2), Inches(1.55), 5.2)
caption(s, "생육·수확량 예측 화면", Inches(10.2), Inches(7.0))

# ───────────────────────── 보안·운영 ─────────────────────────
s = prs.slides.add_slide(BLANK); _bar(s, "보안 · 운영", "Security & Ops")
bullets(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(2.6),
    ["무인증 데이터 수집·모델오염·재학습 DoS 차단 (인증+소유권)",
     "WebSocket/센서 무인증 도청 차단 (토큰·소유권 검증)",
     "공개 데모 무자격 토큰 전환 — 소스 내 자격증명 제거",
     "fail-closed 인증 · 농가 격리 · 결제·연합 무결성 · 클러스터 PII 익명화"],
    size=14.5, head="외부 보안감사 P0~P3 전 항목 조치 완료")
bullets(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.6),
    ["Cloudflare named tunnel → uvicorn :8000 (HTTPS)",
     "watchdog 30초 자동복구 + 단일 인스턴스 가드",
     "PUBLIC_DEMO 읽기전용 게이트(파괴·고비용 작업 차단)",
     "PWA 서비스워커 캐시 버저닝 · SEO(robots·sitemap·구조화데이터)"],
    size=14.5, head="운영 인프라")

# ───────────────────────── 마무리 ─────────────────────────
s = prs.slides.add_slide(BLANK)
_box(s, 0, 0, SW, SH, fill=GREEN_DARK)
_text(s, Inches(1), Inches(2.7), Inches(11.3), Inches(2.0),
      [[("KAASA Farmingsight", 40, True, WHITE)],
       [("데이터로 농사를 결정하다 — 환경·생육·경영을 하나로", 18, False, RGBColor(0xCF,0xE9,0xDB))]],
      align=PP_ALIGN.CENTER)
_text(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.6),
      [[("https://farmingsight.org", 16, True, GREEN)]], align=PP_ALIGN.CENTER)

prs.save(OUT)
print("저장 완료:", OUT, "| 슬라이드", len(prs.slides._sldIdLst))
