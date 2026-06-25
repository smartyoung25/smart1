# -*- coding: utf-8 -*-
"""KAASA smartfarmingsight 통합 소개서 (PPTX, 16:9) — 시스템 소개 + 7단계 로드맵, 디자인 통일."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

SHOTS = os.path.join("out", "ppt_shots")
OUT = os.path.join("out", "KAASA_smartfarmingsight_통합소개서.pptx")

GREEN_DARK = RGBColor(0x0F, 0x51, 0x32)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
INK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x5A, 0x6A, 0x60)
LIGHT = RGBColor(0xF2, 0xF6, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
PR = 390 / 844  # phone w/h


def _box(s, x, y, w, h, fill=None, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    return sp


def _text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sa=6):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sa); p.space_before = Pt(0)
        if isinstance(para, tuple): para = [para]
        for (t, sz, b, c) in para:
            r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.bold = b
            r.font.color.rgb = c; r.font.name = FONT
    return tb


def _bar(s, title, kicker=None):
    _box(s, 0, 0, SW, Inches(0.95), fill=GREEN_DARK)
    _box(s, 0, Inches(0.95), SW, Inches(0.06), fill=GREEN)
    _text(s, Inches(0.6), 0, Inches(10.5), Inches(0.95), [[(title, 24, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        _text(s, Inches(9.3), 0, Inches(3.4), Inches(0.95), [[(kicker, 12, False, RGBColor(0xCF,0xE9,0xDB))]],
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def phone(s, name, cx, cy, h_in):
    h = Inches(h_in); w = Emu(int(h * PR)); x = cx - w // 2; pad = Inches(0.05)
    _box(s, x - pad, cy - pad, w + pad*2, h + pad*2, fill=WHITE, line=RGBColor(0xD8,0xE2,0xDC))
    p = os.path.join(SHOTS, name + ".png")
    if os.path.exists(p): s.shapes.add_picture(p, x, cy, height=h)
    return w


def caption(s, text, cx, y, w_in=2.6):
    _text(s, cx - Inches(w_in/2), y, Inches(w_in), Inches(0.4), [[(text, 12, True, GREEN_DARK)]], align=PP_ALIGN.CENTER)


def bullets(s, x, y, w, h, items, size=14, head=None):
    runs = []
    if head: runs.append([(head, 17, True, GREEN_DARK)])
    for it in items: runs.append([("•  ", size, True, GREEN), (it, size, False, INK)])
    _text(s, x, y, w, h, runs, sa=9)


def feature_slide(title, kicker, head, items, shots):
    s = prs.slides.add_slide(BLANK); _bar(s, title, kicker)
    bullets(s, Inches(0.6), Inches(1.45), Inches(5.1), Inches(5.4), items, size=14, head=head)
    for i, (nm, cap) in enumerate(shots[:2]):
        cx = Inches(8.0) if i == 0 else Inches(11.0)
        phone(s, nm, cx, Inches(1.55), 5.0); caption(s, cap, cx, Inches(6.95))
    return s


def divider(title, subtitle):
    s = prs.slides.add_slide(BLANK)
    _box(s, 0, 0, SW, SH, fill=GREEN_DARK)
    _box(s, Inches(0.9), Inches(3.0), Inches(0.18), Inches(1.5), fill=GREEN)
    _text(s, Inches(1.3), Inches(3.0), Inches(11), Inches(1.0), [[(title, 36, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    _text(s, Inches(1.32), Inches(4.15), Inches(11), Inches(0.6), [[(subtitle, 16, False, RGBColor(0xCF,0xE9,0xDB))]])
    return s


# ───────── 1. 타이틀 ─────────
s = prs.slides.add_slide(BLANK); _box(s, 0, 0, SW, SH, fill=GREEN_DARK)
if os.path.exists("og-image.png"):
    bw = Inches(8.2); s.shapes.add_picture("og-image.png", (SW-bw)//2, Inches(0.7), width=bw)
_text(s, Inches(1), Inches(5.0), Inches(11.3), Inches(1.2),
      [[("스마트농업 경영최적화 의사결정 플랫폼", 30, True, WHITE)],
       [("환경관리 · 생육모델 · 이기종 통합 · 온실/노지 관리 — 시스템 소개 & 실증 로드맵", 15, False, RGBColor(0xCF,0xE9,0xDB))]],
      align=PP_ALIGN.CENTER)
_text(s, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5), [[("https://farmingsight.org", 13, False, GREEN)]], align=PP_ALIGN.CENTER)

# ───────── 2. 개요 ─────────
s = prs.slides.add_slide(BLANK); _bar(s, "시스템 개요", "What it is")
_text(s, Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.9),
      [[("환경·기상·시장 데이터를 받아 ", 16, False, INK), ("수확량 → 매출 → 비용 → 순이익", 16, True, GREEN_DARK),
        ("을 예측하고, 최적 환경과 관수·출하 의사결정을 제시하는 스마트팜 운영 OS", 16, False, INK)]])
cards = [("🌡 스마트팜 환경관리","온도·습도·CO₂·VPD·DLI 광연동 전략표·AI 처방"),
         ("💧 일사적산 관수","P1~P6 일일 곡선(J/cm²)·EC·배액률·함수율"),
         ("📈 생육모델·수확량 예측","M1~M5 모델 체인 수확량·순이익 예측"),
         ("🔌 이기종 장비 통합","서로 다른 제조사 제어기·센서를 하나로"),
         ("🏡 온실6+노지7 = 13작목","온실 6종·제주 노지 7종 광역 지원"),
         ("🤝 공동출하·경영분석","시세·채널 비교·월간 경영성과 리포트")]
cx0, cy0, cw, ch, gx, gy = Inches(0.6), Inches(2.25), Inches(3.9), Inches(1.45), Inches(0.18), Inches(0.22)
for i,(t,d) in enumerate(cards):
    r,c = divmod(i,3); x = cx0 + c*(cw+gx); y = cy0 + r*(ch+gy)
    _box(s,x,y,cw,ch,fill=LIGHT,line=RGBColor(0xDD,0xE7,0xE1))
    _text(s,x+Inches(0.18),y+Inches(0.12),cw-Inches(0.36),ch-Inches(0.24),[[(t,14,True,GREEN_DARK)],[(d,11.5,False,GRAY)]],sa=4)

# ───────── 3. 아키텍처 ─────────
s = prs.slides.add_slide(BLANK); _bar(s, "핵심 아키텍처", "Architecture")
layers = [("프런트엔드","모바일 화면 41종 · 공용 레이어 · PWA(오프라인·캐시)", GREEN),
          ("API (FastAPI)","JWT 인증 · 농가 소유권 · PUBLIC_DEMO 게이트 · 라우터(farmer·admin·billing·ws…)", RGBColor(0x2D,0x9C,0xDB)),
          ("분석 코어","M1 생육 → M2 수확량 → M3 매출 → M4 비용 → M5 병해 · profit_optimizer", RGBColor(0xE6,0x7E,0x22)),
          ("파이프라인","ETL 증분 · 임계 재학습·배포 게이트 · 모델 레지스트리 · 연합학습 · MQTT", RGBColor(0x8E,0x44,0xAD)),
          ("배포·운영","Cloudflare tunnel → uvicorn :8000 · watchdog 자동복구 · farmingsight.org", GREEN_DARK)]
y = Inches(1.4)
for title,desc,col in layers:
    _box(s,Inches(0.6),y,Inches(12.1),Inches(0.98),fill=LIGHT,line=RGBColor(0xDD,0xE7,0xE1))
    _box(s,Inches(0.6),y,Inches(0.16),Inches(0.98),fill=col)
    _text(s,Inches(1.0),y,Inches(3.0),Inches(0.98),[[(title,16,True,col)]],anchor=MSO_ANCHOR.MIDDLE)
    _text(s,Inches(4.0),y,Inches(8.5),Inches(0.98),[[(desc,12.5,False,INK)]],anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.12)

# ───────── 4~11. 기능 화면 ─────────
feature_slide("랜딩 · 화면 네비게이터","Onboarding","첫 진입 & 탐색",
  ["키워드 중심 랜딩(SEO 최적화)으로 가치 전달","화면 네비게이터에서 41개 화면 탐색",
   "등급별 맞춤 메뉴·잠금 배지로 접근 경로 안내","PWA 설치·오프라인 캐시"],
  [("intro","랜딩(소개)"),("smartos","화면 네비게이터")])
feature_slide("통합 홈 · 온실 홈","Dashboard","일일 운영 대시보드",
  ["VPD·관수·생육·병해 핵심 지표를 한 화면에","의사결정 카드 + '적용 기록' 폐루프",
   "실시간 센서 WebSocket(인증·소유권 검증)","작물 전환으로 13작목 즉시 비교"],
  [("c3_home","통합 홈"),("g1_home","온실 홈")])
feature_slide("환경관리 전략 · 관수 P1~P6","Climate & Irrigation","정밀 제어",
  ["생육시기×하루4구간 환경 전략표(온·습도·CO₂·VPD)","목표 대비 실측 편차 → AI 처방",
   "관수는 일사 적산(J/cm²) 기준 P1~P6","급액/배액 EC·pH·배액률·함수율 관리"],
  [("g2_env","환경관리 전략"),("g3_period","관수 P1~P6")])
feature_slide("생육·수확량 예측 · 병해 조기경보","ML & IPM","예측과 예방",
  ["M1~M2 생육모델로 수확량(kg/m²) 예측","예측 vs 실측 피드백 → 임계 재학습",
   "결로시간 기반 병해 조기경보(IPM)","방제 이행 기록으로 학습·품질 폐루프"],
  [("g4_growth","생육·수확량 예측"),("g5_disease","병해 조기경보")])
feature_slide("노지 관리 · 경영·ERP 분석","Field & ERP","노지와 경영",
  ["제주 노지 7작목 — 토양수분·관개·기상","경영 손익(매출·비용·순이익) 분해",
   "절감 조치 이행 기록으로 비용 최적화","현장 데이터 입력 → 모델 학습 피드"],
  [("f1_field","노지 관리"),("c5_erp","경영·ERP 분석")])
feature_slide("공동출하 · 월간 경영성과 리포트","Market & Report","유통과 성과",
  ["시세·유통 채널 비교, 공동출하 참여","작목별 등급 자동 적용(번과/특·상·등외)",
   "월간 경영성과 리포트(정책 지표 연동)","성과지표 변화율·월 스냅샷 이력"],
  [("c12_joint","공동출하"),("c14_report","월간 경영성과 리포트")])
feature_slide("등급 비교 · 클러스터 관제","Tiers & Cluster","구독과 광역 관제",
  ["4등급 기능 매트릭스·AI 쿼터 가시화","업그레이드 안내",
   "다중농가 클러스터 광역 작황·진단 집계","이상 농가 위치특정(공개 응답 PII 익명화)"],
  [("c22_tiers","등급 비교"),("c20_cluster","클러스터 관제")])
feature_slide("AI 영농비서 · 이기종 장비 통합","AI & Interop","대화형 운영과 통합",
  ["내 농장 데이터 기반 AI 영농비서(RAG)","관수·환경·수확·병해 질의응답",
   "서로 다른 제조사 제어기·센서 등록","데이터 통합·연동 신청으로 단일 화면"],
  [("c13_chat","AI 영농비서"),("c8_interop","이기종 장비 통합")])

# ───────── 12. 작물·모델 + 성능 ─────────
s = prs.slides.add_slide(BLANK); _bar(s, "작물 커버리지 & 모델 성능", "Crops & Models")
bullets(s, Inches(0.6), Inches(1.35), Inches(6.2), Inches(5.4),
  ["온실 6작목(딸기·오이·완숙·방울토마토·파프리카·참외)","제주 노지 7작목 — 총 13작목",
   "M1 생육 · M2 수확량 · M3 매출 · M4 비용 · M5 병해","수확량 예측 MAPE 8.7~19.8%(게이트 35% 이내)",
   "생육 설명력 R² 0.17~0.37 — 데이터 축적 시 개선","드리프트 모니터링 + 임계 자동 재학습·폴백"],
  size=14.5, head="13작목 × M1~M5 모델 체인")
phone(s, "g4_growth", Inches(10.2), Inches(1.55), 5.2); caption(s, "생육·수확량 예측", Inches(10.2), Inches(7.0))

# ───────── 13. 보안·운영 ─────────
s = prs.slides.add_slide(BLANK); _bar(s, "보안 · 운영", "Security & Ops")
bullets(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(2.6),
  ["무인증 데이터 수집·모델오염·재학습 DoS 차단(인증+소유권)","WebSocket/센서 무인증 도청 차단",
   "공개 데모 무자격 토큰 — 소스 내 자격증명 제거","fail-closed 인증·농가 격리·결제/연합 무결성·클러스터 PII 익명화"],
  size=14.5, head="외부 보안감사 P0~P3 전 항목 조치 완료")
bullets(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.6),
  ["Cloudflare tunnel → uvicorn :8000 (HTTPS)","watchdog 30초 자동복구 + 단일 인스턴스 가드",
   "PUBLIC_DEMO 읽기전용 게이트(파괴·고비용 작업 차단)","PWA 캐시 버저닝 · SEO(robots·sitemap·구조화데이터)"],
  size=14.5, head="운영 인프라")

# ═════════ 14. 로드맵 섹션 divider ═════════
divider("7단계 실증 추진 로드맵", "참여농가 사전진단부터 데이터 수집·AI 의사결정·컨설팅·사업화 모델 도출까지")

# ───────── 15. 로드맵 요약 ─────────
s = prs.slides.add_slide(BLANK); _bar(s, "실증 추진 7단계 — 한눈에", "Roadmap")
steps_sum = [("1","참여농가 사전진단 & 현장 문제 유형화","C17·C18"),
             ("2","데이터 수집항목·주기·입력방식 확정","C21·C2"),
             ("3","농가별 데이터 수집 & 품질관리","C16·C8"),
             ("4","AI 의사결정 지원 대시보드 시범적용","C3·G2"),
             ("5","농가별 분석리포트 & 현장 컨설팅","C14·C4"),
             ("6","농가 피드백 수렴 & 서비스 개선","C13·도움말"),
             ("7","실증성과 분석 & 사업화 모델 도출","C9·C10")]
y = Inches(1.3)
for n, t, sc in steps_sum:
    _box(s, Inches(0.6), y, Inches(0.62), Inches(0.62), fill=GREEN)
    _text(s, Inches(0.6), y, Inches(0.62), Inches(0.62), [[(n, 20, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _box(s, Inches(1.4), y, Inches(9.0), Inches(0.62), fill=LIGHT, line=RGBColor(0xDD,0xE7,0xE1))
    _text(s, Inches(1.65), y, Inches(8.6), Inches(0.62), [[(t, 14, True, INK)]], anchor=MSO_ANCHOR.MIDDLE)
    _box(s, Inches(10.6), y, Inches(2.1), Inches(0.62), fill=GREEN_DARK)
    _text(s, Inches(10.6), y, Inches(2.1), Inches(0.62), [[(sc, 12, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.78)

# ───────── 16~22. 7단계 상세(스크린샷) ─────────
ROAD = [
  ("1단계 · 사전진단 & 현장 문제 유형화","Diagnosis",
   ["[C17] 성숙도 종합점수·6대영역 진단·우선순위 처방·ROI·결과 PDF",
    "[C18] RDA 현장 체크리스트로 센서 밖 항목(원수·필터·훈증·병해충·인증)까지 반영해 문제 유형화"],
   [("c17_diagnosis","C17 시스템 종합진단"),("c18_checklist","C18 현장 컨설팅 문진")]),
  ("2단계 · 수집항목·주기·입력방식 확정","Data Scope",
   ["[C21] 장비·서비스·프로토콜·인프라정도·요청내용으로 이기종·외부데이터·전문가 연동 신청·관리",
    "[C2] 활용 모드·항목별 공개 토글(민감정보 비공개, 학습·벤치마킹만)로 수집·활용범위 확정"],
   [("c21_apply","C21 연동·서비스 신청"),("c2_consent","C2 데이터 동의")]),
  ("3단계 · 데이터 수집 & 품질관리","Ingest & QA",
   ["[C16] 이기종 장비 등록·표준변수 매핑(견적서·시방서·PDF 자동추출), 분류·ID·주소·포트·프로토콜",
    "[C8] 연동API·연결·연동률·외부데이터 실시간·MQTT·게이트웨이로 품질·연동상태 관리"],
   [("c16_equipment","C16 시설기자재 등록"),("c8_interop","C8 이기종 연동")]),
  ("4단계 · AI 의사결정 대시보드 시범적용","Decision Support",
   ["[C3] 오늘의 결정·수확예측·예상매출·배액률·VPD·벤치마크·AI 추천(관수 신뢰도) 통합 표시",
    "[G2] 온·습도·CO₂·VPD·DLI 실측·전략표·제어모드 4단계·AI 에이전트(30초)로 시범적용"],
   [("c3_home","C3 통합 홈"),("g2_env","G2 환경·기후 제어")]),
  ("5단계 · 분석리포트 & 현장 컨설팅","Report & Consulting",
   ["[C14] 6대영역·전월대비·9대 성과지표·이행활동·취약항목·실행 To-do·PDF",
    "[C4] AI 종합진단·1차 개선추천·환경이상 감지·전문가 컨설팅 예약(현장결과 자동첨부)"],
   [("c14_report","C14 월간 경영성과 리포트"),("c4_diagnosis","C4 AI 진단")]),
  ("6단계 · 피드백 수렴 & 기능 개선","Feedback",
   ["[C13] 농장데이터+농진청·병해충DB(RAG) 결합 생성형 AI 챗봇 Q&A·피드백 수렴(출처배지)",
    "[도움말] 영농가이드·빠른시작·FAQ·용어·출처배지 매뉴얼로 사용성 개선"],
   [("c13_chat","C13 AI 영농비서"),("help","도움말·사용자 매뉴얼")]),
  ("7단계 · 실증성과 분석 & 사업화","Outcome & Scale",
   ["[C9] RDA 표준·우수농가 대비 비교·개선포인트·월간리포트 연계로 성과 분석",
    "[C10] 연매출·연에너지비 기반 ROI·투자안 비교·회수기간으로 사업화·확산 전략 도출"],
   [("c9_benchmark","C9 벤치마킹"),("c10_roi","C10 투자 ROI")]),
]
for title, kicker, items, shots in ROAD:
    feature_slide(title, kicker, "추진 내용", items, shots)

# ───────── 마무리 ─────────
s = prs.slides.add_slide(BLANK); _box(s, 0, 0, SW, SH, fill=GREEN_DARK)
_text(s, Inches(1), Inches(2.7), Inches(11.3), Inches(2.0),
      [[("KAASA smartfarmingsight", 40, True, WHITE)],
       [("데이터로 농사를 결정하다 — 환경·생육·경영을 하나로", 18, False, RGBColor(0xCF,0xE9,0xDB))]], align=PP_ALIGN.CENTER)
_text(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.6), [[("https://farmingsight.org", 16, True, GREEN)]], align=PP_ALIGN.CENTER)

prs.save(OUT)
print("저장 완료:", OUT, "| 슬라이드", len(prs.slides._sldIdLst))
