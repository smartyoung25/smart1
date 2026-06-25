# -*- coding: utf-8 -*-
"""KAASA smartfarmingsight 도입 제안서 (PPTX, 16:9) — 회원가입 유도(CTA)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

SHOTS = os.path.join("out", "ppt_shots")
OUT = os.path.join("out", "KAASA_smartfarmingsight_제안서.pptx")

GREEN_DARK = RGBColor(0x0F, 0x51, 0x32)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
INK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x5A, 0x6A, 0x60)
LIGHT = RGBColor(0xF2, 0xF6, 0xF4)
RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]; PR = 390/844


def _box(s,x,y,w,h,fill=None,line=None,lw=1):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    return sp

def _round(s,x,y,w,h,fill,line=None):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h); sp.shadow.inherit=False
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1.5)
    return sp

def _text(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sa=6):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sa); p.space_before=Pt(0)
        if isinstance(para,tuple): para=[para]
        for (t,sz,b,c) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=FONT
    return tb

def _bar(s,title,kicker=None):
    _box(s,0,0,SW,Inches(0.95),fill=GREEN_DARK); _box(s,0,Inches(0.95),SW,Inches(0.06),fill=GREEN)
    _text(s,Inches(0.6),0,Inches(10.5),Inches(0.95),[[(title,24,True,WHITE)]],anchor=MSO_ANCHOR.MIDDLE)
    if kicker: _text(s,Inches(9.3),0,Inches(3.4),Inches(0.95),[[(kicker,12,False,RGBColor(0xCF,0xE9,0xDB))]],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)

def phone(s,name,cx,cy,h_in):
    h=Inches(h_in); w=Emu(int(h*PR)); x=cx-w//2; pad=Inches(0.05)
    _box(s,x-pad,cy-pad,w+pad*2,h+pad*2,fill=WHITE,line=RGBColor(0xD8,0xE2,0xDC))
    p=os.path.join(SHOTS,name+".png")
    if os.path.exists(p): s.shapes.add_picture(p,x,cy,height=h)
    return w

def caption(s,t,cx,y,w=2.6):
    _text(s,cx-Inches(w/2),y,Inches(w),Inches(0.4),[[(t,12,True,GREEN_DARK)]],align=PP_ALIGN.CENTER)

def cards4(s, items, y0=Inches(1.45)):
    cw,ch,gx,gy=Inches(5.85),Inches(2.15),Inches(0.3),Inches(0.3)
    x0=Inches(0.6)
    for i,(ico,t,d,col) in enumerate(items):
        r,c=divmod(i,2); x=x0+c*(cw+gx); y=y0+r*(ch+gy)
        _round(s,x,y,cw,ch,LIGHT,line=RGBColor(0xDD,0xE7,0xE1))
        _box(s,x,y,Inches(0.14),ch,fill=col)
        _text(s,x+Inches(0.35),y+Inches(0.22),cw-Inches(0.6),Inches(0.6),[[(ico+"  "+t,17,True,col)]])
        _text(s,x+Inches(0.35),y+Inches(0.92),cw-Inches(0.6),ch-Inches(1.1),[[(d,13.5,False,INK)]])

# ───────── 1. 표지(제안서) ─────────
s=prs.slides.add_slide(BLANK); _box(s,0,0,SW,SH,fill=GREEN_DARK)
if os.path.exists("og-image.png"):
    bw=Inches(7.4); s.shapes.add_picture("og-image.png",(SW-bw)//2,Inches(0.6),width=bw)
_text(s,Inches(1),Inches(4.7),Inches(11.3),Inches(1.4),
      [[("스마트팜 경영최적화 도입 제안", 32, True, WHITE)],
       [("데이터로 농사를 결정하다 — 환경·생육·경영을 하나로", 16, False, RGBColor(0xCF,0xE9,0xDB))]],align=PP_ALIGN.CENTER)
_round(s,Inches(5.0),Inches(6.35),Inches(3.33),Inches(0.6),GREEN)
_text(s,Inches(5.0),Inches(6.35),Inches(3.33),Inches(0.6),[[("무료로 시작하기 →  farmingsight.org",13,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# ───────── 2. 농가의 고민(문제) ─────────
s=prs.slides.add_slide(BLANK); _bar(s,"이런 고민, 있으셨나요?","Problem")
cards4(s,[
  ("🤔","감과 경험에 의존","환경·관수·시비 판단이 데이터 없이 직관에 의존 → 작황 편차가 크고 재현이 어렵습니다.",RED),
  ("🧩","데이터가 흩어져 있음","센서·기상·시장·장비 데이터가 단절돼 경영 의사결정에 활용하지 못합니다.",ORANGE),
  ("🔌","이기종 장비 파편화","제조사마다 제어기·앱이 달라 통합 관제·비교가 불가능합니다.",RED),
  ("💸","경영이 깜깜이","수확량·매출·비용·순이익을 미리 못 봐 투자·출하 판단이 어렵습니다.",ORANGE)])

# ───────── 3. 해결(가치 제안) ─────────
s=prs.slides.add_slide(BLANK); _bar(s,"KAASA가 해결해 드립니다","Solution")
cards4(s,[
  ("📊","데이터 기반 의사결정","환경·관수·시비를 일사적산·VPD·전략표로 정밀 처방 — 감 대신 데이터로.",GREEN_DARK),
  ("🗂","흩어진 데이터를 하나로","센서·기상·시장·장비를 표준 변수로 통합해 한 화면에서 운영.",GREEN_DARK),
  ("🔗","이기종 장비 표준 통합","서로 다른 제조사 제어기·센서를 MQTT·게이트웨이로 단일 관제.",GREEN_DARK),
  ("📈","수확량·순이익 예측 경영","M1~M5 모델로 수확량→매출→비용→순이익을 미리 예측·최적화.",GREEN_DARK)])

# ───────── 4~6. 핵심 기능 미리보기 ─────────
def feat(title,kicker,head,items,shots):
    s=prs.slides.add_slide(BLANK); _bar(s,title,kicker)
    runs=[[(head,17,True,GREEN_DARK)]]+[[("•  ",14,True,GREEN),(it,14,False,INK)] for it in items]
    _text(s,Inches(0.6),Inches(1.45),Inches(5.1),Inches(5.4),runs,sa=10)
    for i,(nm,cap) in enumerate(shots[:2]):
        cx=Inches(8.0) if i==0 else Inches(11.0)
        phone(s,nm,cx,Inches(1.55),5.0); caption(s,cap,cx,Inches(6.95))
feat("환경관리 · 일사적산 관수","Feature 1","현장을 정밀 제어",
  ["온·습도·CO₂·VPD·DLI 실측 + 환경 전략표","목표 대비 편차 → AI 처방(제어 4단계)",
   "시각이 아닌 일사 적산(J/cm²) 기반 P1~P6 관수","급액/배액 EC·pH·배액률·함수율 관리"],
  [("g2_env","환경관리 전략"),("g3_period","관수 P1~P6")])
feat("생육·수확량 예측 · 병해 예방","Feature 2","예측과 예방으로 손실↓",
  ["작물별 수확량(kg/m²) 예측 — 출하·매출 계획","예측 vs 실측 피드백 → 모델 자동 개선",
   "결로시간 기반 병해 조기경보(IPM)","방제 이행 기록으로 품질·학습 폐루프"],
  [("g4_growth","수확량 예측"),("g5_disease","병해 조기경보")])
feat("경영 분석 · 월간 성과 리포트","Feature 3","경영을 한눈에",
  ["매출·비용·순이익 분해 + ROI 산출","월간 경영성과 리포트(9대 성과지표)",
   "우수농가 벤치마크로 개선 포인트 제시","AI 영농비서로 즉시 질의응답"],
  [("c14_report","월간 경영성과 리포트"),("c10_roi","투자 ROI")])

# ───────── 7. 왜 KAASA인가(차별점) ─────────
s=prs.slides.add_slide(BLANK); _bar(s,"왜 KAASA인가","Why us")
diffs=[("🌞 일사적산 정밀 관수","시각이 아닌 광량(J/cm²) 기반 — Priva·Grodan 글로벌 기준 정합"),
       ("🧪 13작목 검증 모델","온실 6+노지 7 작목, 수확량 예측 오차(MAPE) 8.7~19.8%"),
       ("🔌 이기종 장비 통합","제조사 무관 표준 변수 매핑·MQTT·게이트웨이 단일 관제"),
       ("🔁 폐루프 학습","현장 이행 데이터가 모델을 지속 개선 — 쓸수록 똑똑해짐"),
       ("🛡 운영급 보안·안정","외부 보안감사 전 항목 조치·HTTPS·PWA·자동복구 watchdog")]
y=Inches(1.4)
for t,d in diffs:
    _round(s,Inches(0.6),y,Inches(12.1),Inches(0.92),LIGHT,line=RGBColor(0xDD,0xE7,0xE1))
    _text(s,Inches(0.95),y,Inches(4.2),Inches(0.92),[[(t,15,True,GREEN_DARK)]],anchor=MSO_ANCHOR.MIDDLE)
    _text(s,Inches(5.2),y,Inches(7.3),Inches(0.92),[[(d,13,False,INK)]],anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(1.05)

# ───────── 8. 요금제(무료 시작) ─────────
s=prs.slides.add_slide(BLANK); _bar(s,"무료로 시작하세요 · 요금제","Pricing")
tiers=[("Basic","무료","핵심 대시보드·기록·기본 분석",GREEN),
       ("Smart","구독","환경 전략표·관수 정밀·AI 추천",RGBColor(0x2D,0x9C,0xDB)),
       ("Pro","구독","수확 예측·경영 리포트·벤치마크",ORANGE),
       ("Enterprise","구독","클러스터 관제·연동·전문 컨설팅",RGBColor(0x8E,0x44,0xAD))]
cw,gx=Inches(2.85),Inches(0.27); x=Inches(0.7)
for nm,price,desc,col in tiers:
    _round(s,x,Inches(1.5),cw,Inches(3.6),WHITE,line=col)
    _box(s,x,Inches(1.5),cw,Inches(0.85),fill=col)
    _text(s,x,Inches(1.5),cw,Inches(0.85),[[(nm,17,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    _text(s,x,Inches(2.5),cw,Inches(0.7),[[(price,22,True,col)]],align=PP_ALIGN.CENTER)
    _text(s,x+Inches(0.25),Inches(3.3),cw-Inches(0.5),Inches(1.6),[[(desc,12.5,False,INK)]],align=PP_ALIGN.CENTER)
    x+=cw+gx
_text(s,Inches(0.7),Inches(5.4),Inches(12),Inches(0.6),
      [[("✓ Basic 무료 플랜으로 즉시 시작 · 신용카드 불필요 · 언제든 업그레이드",15,True,GREEN_DARK)]],align=PP_ALIGN.CENTER)
phone(s,"c22_tiers",Inches(11.4),Inches(5.9),1.4)

# ───────── 9. 참여 절차(7단계 요약) ─────────
s=prs.slides.add_slide(BLANK); _bar(s,"참여 절차 — 실증 추진 7단계","Process")
steps=[("1","사전진단·문제 유형화"),("2","수집항목·입력방식 확정"),("3","데이터 수집·품질관리"),
       ("4","AI 의사결정 시범적용"),("5","분석리포트·현장 컨설팅"),("6","피드백 수렴·개선"),("7","성과 분석·사업화")]
y=Inches(1.5)
for n,t in steps:
    _box(s,Inches(1.2),y,Inches(0.6),Inches(0.6),fill=GREEN)
    _text(s,Inches(1.2),y,Inches(0.6),Inches(0.6),[[(n,18,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    _round(s,Inches(2.0),y,Inches(9.8),Inches(0.6),LIGHT,line=RGBColor(0xDD,0xE7,0xE1))
    _text(s,Inches(2.3),y,Inches(9.3),Inches(0.6),[[(t,15,True,INK)]],anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(0.74)

# ───────── 10. ★ 회원가입 CTA ─────────
s=prs.slides.add_slide(BLANK); _box(s,0,0,SW,SH,fill=GREEN_DARK)
_text(s,Inches(0.8),Inches(0.7),Inches(8.0),Inches(1.2),
      [[("지금, 무료로 시작하세요",34,True,WHITE)],
       [("회원가입 후 농장을 등록하면 AI 진단을 무료로 체험할 수 있습니다.",16,False,RGBColor(0xCF,0xE9,0xDB))]])
steps=[("1","farmingsight.org 접속","웹·모바일 어디서나, 설치 없이"),
       ("2","회원가입(이메일·전화)","역할 선택 — 농가·조합·유통·전문가·공공"),
       ("3","농장 세팅","지역·작목·재배방식·장비 등록"),
       ("4","AI 진단 무료 체험","시스템 종합진단·환경/관수 처방 즉시 확인")]
y=Inches(2.1)
for n,t,d in steps:
    _box(s,Inches(0.8),y,Inches(0.55),Inches(0.55),fill=GREEN)
    _text(s,Inches(0.8),y,Inches(0.55),Inches(0.55),[[(n,17,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    _text(s,Inches(1.55),y-Inches(0.05),Inches(6.6),Inches(0.7),
          [[(t,16,True,WHITE)],[(d,12,False,RGBColor(0xCF,0xE9,0xDB))]],sa=0)
    y+=Inches(0.95)
# QR 카드
_box(s,Inches(8.7),Inches(1.9),Inches(3.9),Inches(4.7),fill=WHITE)
if os.path.exists("out/qr_farmingsight.png"):
    s.shapes.add_picture("out/qr_farmingsight.png",Inches(9.45),Inches(2.2),height=Inches(2.3))
_text(s,Inches(8.7),Inches(4.65),Inches(3.9),Inches(0.4),[[("QR 스캔 → 바로 접속",13,True,GREEN_DARK)]],align=PP_ALIGN.CENTER)
_text(s,Inches(8.7),Inches(5.05),Inches(3.9),Inches(0.4),[[("farmingsight.org",17,True,GREEN_DARK)]],align=PP_ALIGN.CENTER)
_text(s,Inches(8.7),Inches(5.6),Inches(3.9),Inches(0.9),
      [[("✓ 무료 플랜 · 신용카드 불필요",12.5,True,GREEN_DARK)],
       [("✓ 13작목 · 온실/노지 지원",12.5,True,GREEN_DARK)]],align=PP_ALIGN.CENTER,sa=3)

# ───────── 11. 마무리 ─────────
s=prs.slides.add_slide(BLANK); _box(s,0,0,SW,SH,fill=GREEN_DARK)
_text(s,Inches(1),Inches(2.6),Inches(11.3),Inches(2.0),
      [[("함께 시작하시죠",38,True,WHITE)],
       [("KAASA smartfarmingsight · 데이터로 농사를 결정하다",17,False,RGBColor(0xCF,0xE9,0xDB))]],align=PP_ALIGN.CENTER)
_round(s,Inches(4.7),Inches(4.7),Inches(3.93),Inches(0.7),GREEN)
_text(s,Inches(4.7),Inches(4.7),Inches(3.93),Inches(0.7),[[("무료 회원가입 →  farmingsight.org",14,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print("저장 완료:", OUT, "| 슬라이드", len(prs.slides._sldIdLst))
